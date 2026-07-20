"""Unit tests for PDF and PPTX renderers given a fixture BriefDocument."""

from __future__ import annotations

from datetime import UTC, datetime

from app.domain.entities.brief import BriefDocument, BriefStoryItem


def _make_brief(n_stories: int = 2) -> BriefDocument:
    return BriefDocument(
        workspace_id="ws1",
        entity_keyword="Tesla",
        generated_at=datetime(2025, 7, 15, 9, 0, tzinfo=UTC),
        window_days=7,
        subtitle="2 stories",
        stories=[
            BriefStoryItem(
                story_id=f"s{i}",
                title=f'Story {i}: <bold> & "quotes"',
                entity_name="Tesla",
                risk_level=["low", "medium", "high", "critical"][i % 4],
                sentiment_label="neutral",
                article_count=3 + i,
                what_happened=f"What happened {i}.",
                why_happened=f"Why it happened {i}.",
                top_recommendation=f"Consider action {i}.",
                top_rec_confidence=0.8 + i * 0.05,
                published_at="2025-07-14T10:00:00+00:00",
            )
            for i in range(n_stories)
        ],
    )


# ---------------------------------------------------------------------------
# PDF renderer
# ---------------------------------------------------------------------------


class TestPdfRenderer:
    def test_render_pdf_returns_bytes(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        brief = _make_brief()
        result = render_pdf(brief)
        assert isinstance(result, bytes)

    def test_render_pdf_is_non_empty(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        result = render_pdf(_make_brief())
        assert len(result) > 0

    def test_render_pdf_starts_with_pdf_magic(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        result = render_pdf(_make_brief())
        assert result[:4] == b"%PDF"

    def test_render_pdf_empty_brief_no_crash(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        brief = _make_brief(0)
        result = render_pdf(brief)
        assert isinstance(result, bytes) and len(result) > 0

    def test_render_pdf_escapes_special_chars(self) -> None:
        """XSS/injection in story titles should not produce malformed HTML."""
        from app.infrastructure.exports.pdf_renderer import _build_html

        brief = _make_brief(1)
        brief.stories[0].title = '<script>alert("xss")</script>'
        html = _build_html(brief)
        assert "<script>" not in html
        assert "&lt;script&gt;" in html

    def test_render_pdf_single_story(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        result = render_pdf(_make_brief(1))
        assert isinstance(result, bytes) and result[:4] == b"%PDF"

    def test_render_pdf_many_stories(self) -> None:
        from app.infrastructure.exports.pdf_renderer import render_pdf

        result = render_pdf(_make_brief(10))
        assert isinstance(result, bytes) and len(result) > 0

    def test_html_contains_entity_keyword(self) -> None:
        from app.infrastructure.exports.pdf_renderer import _build_html

        brief = _make_brief(1)
        html = _build_html(brief)
        assert "Tesla" in html

    def test_html_contains_window_label(self) -> None:
        from app.infrastructure.exports.pdf_renderer import _build_html

        brief = _make_brief(1)
        html = _build_html(brief)
        assert "7" in html  # window_days

    def test_html_contains_generated_date(self) -> None:
        from app.infrastructure.exports.pdf_renderer import _build_html

        brief = _make_brief(1)
        html = _build_html(brief)
        assert "2025" in html


# ---------------------------------------------------------------------------
# PPTX renderer
# ---------------------------------------------------------------------------


class TestPptxRenderer:
    def test_render_pptx_returns_bytes(self) -> None:
        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief())
        assert isinstance(result, bytes)

    def test_render_pptx_is_non_empty(self) -> None:
        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief())
        assert len(result) > 0

    def test_render_pptx_starts_with_zip_magic(self) -> None:
        """PPTX files are ZIP archives — start with PK (0x504B)."""
        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief())
        assert result[:2] == b"PK"

    def test_render_pptx_empty_brief_no_crash(self) -> None:
        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief(0))
        assert isinstance(result, bytes) and len(result) > 0

    def test_render_pptx_slide_count(self) -> None:
        """cover + N story slides + closing = N + 2 slides."""
        from io import BytesIO

        from pptx import Presentation

        from app.infrastructure.exports.pptx_renderer import render_pptx

        n = 3
        result = render_pptx(_make_brief(n))
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == n + 2  # cover + stories + closing

    def test_render_pptx_single_story_slide_count(self) -> None:
        from io import BytesIO

        from pptx import Presentation

        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief(1))
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 3  # cover + 1 + closing

    def test_render_pptx_many_stories(self) -> None:
        from app.infrastructure.exports.pptx_renderer import render_pptx

        result = render_pptx(_make_brief(15))
        assert isinstance(result, bytes) and result[:2] == b"PK"

    def test_render_pptx_no_recommendation_no_crash(self) -> None:
        """Story with empty top_recommendation should still render."""
        from app.infrastructure.exports.pptx_renderer import render_pptx

        brief = _make_brief(1)
        brief.stories[0].top_recommendation = ""
        result = render_pptx(brief)
        assert isinstance(result, bytes) and len(result) > 0
