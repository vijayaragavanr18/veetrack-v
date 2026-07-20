"""PPTX renderer — converts a BriefDocument to PPTX bytes via python-pptx.

One slide per story (title + what_happened + why_happened + top recommendation).
Cover slide + closing slide included.
"""

from __future__ import annotations

import io

from app.domain.entities.brief import BriefDocument, BriefStoryItem

_RISK_RGB = {
    "critical": (239, 68, 68),
    "high": (249, 115, 22),
    "medium": (234, 179, 8),
    "low": (34, 197, 94),
}
_DARK = (15, 23, 42)  # #0f172a
_SLATE = (100, 116, 139)  # #64748b
_WHITE = (248, 250, 252)  # #f8fafc
_BODY = (30, 41, 59)  # #1e293b


def render_pptx(brief: BriefDocument) -> bytes:
    """Render *brief* to PPTX bytes.  Raises ImportError if python-pptx is absent."""
    from pptx import Presentation  # deferred import
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    _add_cover_slide(prs, blank_layout, brief)
    for i, story in enumerate(brief.stories):
        _add_story_slide(prs, blank_layout, story, i + 1, len(brief.stories))
    _add_closing_slide(prs, blank_layout)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def _add_cover_slide(prs: object, layout: object, brief: BriefDocument) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)  # type: ignore[attr-defined]
    _fill_bg(slide, _DARK)

    date_str = brief.generated_at.strftime("%B %d, %Y")

    _add_text_box(
        slide,
        "VeeTrack Executive Brief",
        left=Inches(0.8),
        top=Inches(2.4),
        width=Inches(11.7),
        height=Inches(0.9),
        font_size=Pt(36),
        bold=True,
        color=_WHITE,
    )
    _add_text_box(
        slide,
        f"{brief.entity_keyword}  ·  {brief.subtitle}",
        left=Inches(0.8),
        top=Inches(3.4),
        width=Inches(11.7),
        height=Inches(0.5),
        font_size=Pt(18),
        color=(148, 163, 184),
    )
    _add_text_box(
        slide,
        f"Last {brief.window_days} days  ·  Generated {date_str}",
        left=Inches(0.8),
        top=Inches(4.0),
        width=Inches(11.7),
        height=Inches(0.4),
        font_size=Pt(13),
        color=_SLATE,
    )


def _add_story_slide(
    prs: object,
    layout: object,
    story: BriefStoryItem,
    idx: int,
    total: int,
) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)  # type: ignore[attr-defined]
    _fill_bg(slide, _WHITE)

    risk_rgb = _RISK_RGB.get(story.risk_level, _SLATE)

    # Top accent bar
    _add_rect(slide, 0, 0, Inches(13.33), Inches(0.08), risk_rgb)

    # Story number / pagination
    _add_text_box(
        slide,
        f"{idx} / {total}",
        left=Inches(11.8),
        top=Inches(0.15),
        width=Inches(1.4),
        height=Inches(0.3),
        font_size=Pt(10),
        color=_SLATE,
    )

    # Risk badge + entity
    _add_text_box(
        slide,
        f"{story.risk_level.upper()} RISK  ·  {story.entity_name}",
        left=Inches(0.5),
        top=Inches(0.18),
        width=Inches(6),
        height=Inches(0.3),
        font_size=Pt(10),
        bold=True,
        color=risk_rgb,
    )

    # Title
    _add_text_box(
        slide,
        story.title,
        left=Inches(0.5),
        top=Inches(0.55),
        width=Inches(12.3),
        height=Inches(0.8),
        font_size=Pt(22),
        bold=True,
        color=_BODY,
    )

    # What happened
    _add_text_box(
        slide,
        "What Happened",
        left=Inches(0.5),
        top=Inches(1.5),
        width=Inches(5.8),
        height=Inches(0.28),
        font_size=Pt(9),
        bold=True,
        color=_SLATE,
    )
    _add_text_box(
        slide,
        story.what_happened or "Analysis pending.",
        left=Inches(0.5),
        top=Inches(1.82),
        width=Inches(5.8),
        height=Inches(2.1),
        font_size=Pt(12),
        color=_BODY,
        wrap=True,
    )

    # Why it happened
    _add_text_box(
        slide,
        "Why It Happened",
        left=Inches(6.9),
        top=Inches(1.5),
        width=Inches(5.9),
        height=Inches(0.28),
        font_size=Pt(9),
        bold=True,
        color=_SLATE,
    )
    _add_text_box(
        slide,
        story.why_happened or "Analysis pending.",
        left=Inches(6.9),
        top=Inches(1.82),
        width=Inches(5.9),
        height=Inches(2.1),
        font_size=Pt(12),
        color=_BODY,
        wrap=True,
    )

    # Divider line
    _add_rect(slide, Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.01), (226, 232, 240))

    # Recommendation
    if story.top_recommendation:
        _add_text_box(
            slide,
            "Recommended Action",
            left=Inches(0.5),
            top=Inches(4.2),
            width=Inches(12.3),
            height=Inches(0.25),
            font_size=Pt(9),
            bold=True,
            color=_SLATE,
        )
        _add_text_box(
            slide,
            story.top_recommendation,
            left=Inches(0.5),
            top=Inches(4.5),
            width=Inches(12.3),
            height=Inches(0.8),
            font_size=Pt(13),
            color=_BODY,
            wrap=True,
        )


def _add_closing_slide(prs: object, layout: object) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(layout)  # type: ignore[attr-defined]
    _fill_bg(slide, _DARK)
    _add_text_box(
        slide,
        "AI-generated suggestions are advisory only.\nVerify before acting.",
        left=Inches(1.5),
        top=Inches(2.8),
        width=Inches(10.3),
        height=Inches(1.5),
        font_size=Pt(16),
        color=(148, 163, 184),
    )
    _add_text_box(
        slide,
        "VeeTrack",
        left=Inches(1.5),
        top=Inches(4.5),
        width=Inches(4),
        height=Inches(0.4),
        font_size=Pt(14),
        color=_SLATE,
    )


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------


def _fill_bg(slide: object, rgb: tuple[int, int, int]) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill  # type: ignore[attr-defined]
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _add_rect(
    slide: object,
    left: object,
    top: object,
    width: object,
    height: object,
    rgb: tuple[int, int, int],
) -> None:
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(  # type: ignore[attr-defined]
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()


def _add_text_box(
    slide: object,
    text: str,
    *,
    left: object,
    top: object,
    width: object,
    height: object,
    font_size: object,
    color: tuple[int, int, int],
    bold: bool = False,
    wrap: bool = False,
) -> None:
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore[attr-defined]
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
